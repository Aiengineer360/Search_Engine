import os
import pickle
import threading
from collections import defaultdict
from math import log
from tkinter import *
from tkinter import filedialog, messagebox, scrolledtext
from tkinter.ttk import Progressbar, Combobox
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
from exceptions import PendingDeprecationWarning
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import PorterStemmer

# Download NLTK data
nltk.download('punkt')
nltk.download('stopwords')

# Constants
INDEX_FILE = "inverted_index.pkl"
DOCUMENTS_FILE = "documents.pkl"

# Initialize NLTK
stemmer = PorterStemmer()
stop_words = set(stopwords.words('english'))

# Global variables
documents = {}
inverted_index = defaultdict(dict)
document_filenames = {}

# Load saved data if exists
if os.path.exists(INDEX_FILE) and os.path.exists(DOCUMENTS_FILE):
    with open(INDEX_FILE, 'rb') as f:
        inverted_index = pickle.load(f)
    with open(DOCUMENTS_FILE, 'rb') as f:
        documents = pickle.load(f)

# Document Parsing Functions
def read_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ''.join(page.extract_text() for page in reader.pages)
    return text

def read_docx(file_path):
    doc = Document(file_path)
    return '\n'.join(para.text for para in doc.paragraphs)

def read_pptx(file_path):
    ppt = Presentation(file_path)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def read_excel(file_path):
    workbook = openpyxl.load_workbook(file_path)
    text = ""
    for sheet in workbook.sheetnames:
        sheet_obj = workbook[sheet]
        for row in sheet_obj.iter_rows(values_only=True):
            text += ' '.join(str(cell) for cell in row if cell) + '\n'
    return text

def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

# Text Preprocessing
def preprocess_text(text):
    tokens = word_tokenize(text.lower())
    cleaned_tokens = [stemmer.stem(token) for token in tokens if token.isalnum() and token not in stop_words]
    return cleaned_tokens

# Build Inverted Index with TF-IDF
def build_inverted_index(documents):
    inverted_index = defaultdict(dict)
    total_docs = len(documents)
    for doc_id, text in documents.items():
        tokens = preprocess_text(text)
        term_freq = defaultdict(int)
        for token in tokens:
            term_freq[token] += 1
        for token, freq in term_freq.items():
            inverted_index[token][doc_id] = freq
    # Calculate IDF and TF-IDF
    for token in inverted_index:
        idf = log(total_docs / len(inverted_index[token]))
        for doc_id in inverted_index[token]:
            tf = inverted_index[token][doc_id]
            inverted_index[token][doc_id] = tf * idf
    return inverted_index

# Search with TF-IDF Ranking
def search(query, inverted_index, documents):
    query_tokens = preprocess_text(query)
    scores = defaultdict(float)
    for token in query_tokens:
        if token in inverted_index:
            for doc_id, tfidf in inverted_index[token].items():
                scores[doc_id] += tfidf
    return sorted(scores.keys(), key=lambda x: scores[x], reverse=True)

# Save Data to Disk
def save_data():
    with open(INDEX_FILE, 'wb') as f:
        pickle.dump(inverted_index, f)
    with open(DOCUMENTS_FILE, 'wb') as f:
        pickle.dump(documents, f)

# GUI Application
class SearchEngineApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Desktop Search Engine")
        self.root.geometry("800x600")

        # Search Frame
        self.search_frame = Frame(root)
        self.search_frame.pack(pady=10)

        self.label = Label(self.search_frame, text="Enter your search query:")
        self.label.grid(row=0, column=0, padx=5)

        self.entry = Entry(self.search_frame, width=50)
        self.entry.grid(row=0, column=1, padx=5)

        self.search_button = Button(self.search_frame, text="Search", command=self.perform_search)
        self.search_button.grid(row=0, column=2, padx=5)

        # Add Document Frame
        self.add_frame = Frame(root)
        self.add_frame.pack(pady=10)

        self.add_file_button = Button(self.add_frame, text="Add Document", command=self.add_document)
        self.add_file_button.pack(side=LEFT, padx=5)

        self.remove_file_button = Button(self.add_frame, text="Remove Document", command=self.remove_document)
        self.remove_file_button.pack(side=LEFT, padx=5)

        # Results Frame
        self.results_frame = Frame(root)
        self.results_frame.pack(fill=BOTH, expand=True)

        self.results_text = scrolledtext.ScrolledText(self.results_frame, wrap=WORD, width=80, height=20)
        self.results_text.pack(fill=BOTH, expand=True)

        # Pagination Frame
        self.pagination_frame = Frame(root)
        self.pagination_frame.pack(pady=10)

        self.prev_button = Button(self.pagination_frame, text="Previous", command=self.prev_page)
        self.prev_button.pack(side=LEFT, padx=5)

        self.next_button = Button(self.pagination_frame, text="Next", command=self.next_page)
        self.next_button.pack(side=LEFT, padx=5)

        self.page_label = Label(self.pagination_frame, text="Page 1")
        self.page_label.pack(side=LEFT, padx=5)

        # Progress Bar
        self.progress = Progressbar(root, orient=HORIZONTAL, length=400, mode='determinate')
        self.progress.pack(pady=10)

        # Pagination Variables
        self.current_page = 1
        self.results_per_page = 5
        self.search_results = []

    def add_document(self):
        file_path = filedialog.askopenfilename(filetypes=[
            ("PDF Files", "*.pdf"), ("Word Files", "*.docx"), ("PowerPoint Files", "*.pptx"),
            ("Excel Files", "*.xlsx"), ("Text Files", "*.txt")
        ])
        if file_path:
            self.progress['value'] = 0
            self.progress.update()
            threading.Thread(target=self.process_document, args=(file_path,)).start()

    def process_document(self, file_path):
        try:
            if file_path.endswith('.pdf'):
                text = read_pdf(file_path)
            elif file_path.endswith('.docx'):
                text = read_docx(file_path)
            elif file_path.endswith('.pptx'):
                text = read_pptx(file_path)
            elif file_path.endswith('.xlsx'):
                text = read_excel(file_path)
            elif file_path.endswith('.txt'):
                text = read_txt(file_path)
            else:
                messagebox.showerror("Error", "Unsupported file format")
                return
            doc_id = len(documents) + 1
            documents[doc_id] = text
            document_filenames[doc_id] = os.path.basename(file_path)
            global inverted_index
            inverted_index = build_inverted_index(documents)
            save_data()
            self.progress['value'] = 100
            self.progress.update()
            messagebox.showinfo("Success", "Document added successfully")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to process document: {str(e)}")

    def remove_document(self):
        file_path = filedialog.askopenfilename(filetypes=[
            ("PDF Files", "*.pdf"), ("Word Files", "*.docx"), ("PowerPoint Files", "*.pptx"),
            ("Excel Files", "*.xlsx"), ("Text Files", "*.txt")
        ])
        if file_path:
            filename = os.path.basename(file_path)
            doc_id = None
            for id, name in document_filenames.items():
                if name == filename:
                    doc_id = id
                    break
            if doc_id:
                del documents[doc_id]
                del document_filenames[doc_id]
                global inverted_index
                inverted_index = build_inverted_index(documents)
                save_data()
                messagebox.showinfo("Success", "Document removed successfully")
            else:
                messagebox.showerror("Error", "Document not found in index")

    def perform_search(self):
        query = self.entry.get()
        if not query:
            messagebox.showerror("Error", "Please enter a query")
            return
        self.search_results = search(query, inverted_index, documents)
        self.current_page = 1
        self.display_results()

    def display_results(self):
        self.results_text.delete(1.0, END)
        if self.search_results:
            start = (self.current_page - 1) * self.results_per_page
            end = start + self.results_per_page
            for doc_id in self.search_results[start:end]:
                self.results_text.insert(END, f"Document: {document_filenames[doc_id]}\n")
                self.results_text.insert(END, self.highlight_query(self.entry.get(), documents[doc_id][:500]) + "\n\n")
            self.page_label.config(text=f"Page {self.current_page}")
        else:
            self.results_text.insert(END, "No documents found")

    def highlight_query(self, query, text):
        query_tokens = preprocess_text(query)
        highlighted_text = text
        for token in query_tokens:
            highlighted_text = highlighted_text.replace(token, f"**{token}**")
        return highlighted_text

    def prev_page(self):
        if self.current_page > 1:
            self.current_page -= 1
            self.display_results()

    def next_page(self):
        total_pages = (len(self.search_results) + self.results_per_page - 1) // self.results_per_page
        if self.current_page < total_pages:
            self.current_page += 1
            self.display_results()

# Run the Application
if __name__ == "__main__":
    root = Tk()
    app = SearchEngineApp(root)
    root.mainloop()
